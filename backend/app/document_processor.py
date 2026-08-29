"""
Document processing pipeline for StudyMate AI.
Handles: PDF text extraction, image OCR, text chunking, and ChromaDB vector storage.
"""
from __future__ import annotations

import io
import logging
import re
import uuid
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Text Extraction
# ---------------------------------------------------------------------------

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF, including OCR for scanned pages when possible."""
    try:
        import PyPDF2  # type: ignore
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        pages_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages_text.append(text.strip())
        text = "\n\n".join(pages_text)
        if text.strip():
            return text

        # Scanned PDFs have no text layer. PyMuPDF renders their pages for OCR.
        import fitz  # type: ignore
        document = fitz.open(stream=file_bytes, filetype="pdf")
        scanned_pages = []
        for page in document:
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            scanned_pages.append(extract_text_from_image(pixmap.tobytes("png")))
        return "\n\n".join(page for page in scanned_pages if page)
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return ""


def extract_text_from_image(file_bytes: bytes) -> str:
    """Preprocess an image and extract printed or handwritten text with OCR."""
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
        from PIL import ImageEnhance, ImageFilter, ImageOps, ImageSequence  # type: ignore
        try:
            from pillow_heif import register_heif_opener  # type: ignore
            register_heif_opener()
        except ImportError:
            pass
        img = Image.open(io.BytesIO(file_bytes))
        pages = []
        for frame in ImageSequence.Iterator(img):
            frame = ImageOps.exif_transpose(frame).convert("L")
            # Upscaling and contrast normalization materially improve phone scans.
            if max(frame.size) < 2400:
                scale = 2400 / max(frame.size)
                frame = frame.resize((round(frame.width * scale), round(frame.height * scale)))
            frame = ImageEnhance.Contrast(frame).enhance(1.8).filter(ImageFilter.SHARPEN)
            text = pytesseract.image_to_string(frame, config="--psm 6").strip()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except Exception as e:
        logger.warning(f"OCR extraction failed (Tesseract may not be installed): {e}")
        return ""


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Dispatch to the correct extractor based on file type."""
    ft = file_type.lower()
    if ft == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ft in {"image", "scan", "jpg", "jpeg", "png", "tiff", "bmp", "webp"}:
        return extract_text_from_image(file_bytes)
    elif ft in {"txt", "md", "csv", "json", "xml", "html", "htm", "log", "yaml", "yml"}:
        return file_bytes.decode("utf-8", errors="replace")
    elif ft == "docx":
        from docx import Document  # type: ignore
        document = Document(io.BytesIO(file_bytes))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    elif ft == "pptx":
        from pptx import Presentation  # type: ignore
        presentation = Presentation(io.BytesIO(file_bytes))
        return "\n".join(
            shape.text for slide in presentation.slides for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
    elif ft == "xlsx":
        import openpyxl  # type: ignore
        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        return "\n".join(
            " | ".join(str(value) for value in row if value is not None)
            for sheet in workbook.worksheets for row in sheet.iter_rows(values_only=True)
        )
    elif ft == "xls":
        import xlrd  # type: ignore
        workbook = xlrd.open_workbook(file_contents=file_bytes)
        return "\n".join(
            " | ".join(str(value) for value in sheet.row_values(row_index) if value != "")
            for sheet in workbook.sheets() for row_index in range(sheet.nrows)
        )
    else:
        return ""


# ---------------------------------------------------------------------------
# Text Chunking
# ---------------------------------------------------------------------------

def clean_text(text: str) -> str:
    """Normalise whitespace and remove junk characters."""
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^\x20-\x7E\n]", "", text)
    return text.strip()


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 80) -> List[str]:
    """
    Split text into overlapping chunks based on word count.
    Each chunk is ~chunk_size words with `overlap` words carried over.
    """
    words = text.split()
    if not words:
        return []

    chunks: List[str] = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end == len(words):
            break
        start += chunk_size - overlap

    return [c for c in chunks if len(c.strip()) > 30]


# ---------------------------------------------------------------------------
# ChromaDB Vector Storage
# ---------------------------------------------------------------------------

CHROMA_PATH = Path(__file__).parent.parent.parent / "database" / "chroma_store"
CHROMA_PATH.mkdir(parents=True, exist_ok=True)


def _get_chroma_collection(collection_name: str = "studymate_chunks"):
    """Get or create a ChromaDB persistent collection."""
    try:
        import chromadb  # type: ignore
        client = chromadb.PersistentClient(path=str(CHROMA_PATH))
        collection = client.get_or_create_collection(
            name=collection_name,
            metadata={"hnsw:space": "cosine"},
        )
        return collection
    except Exception as e:
        logger.error(f"ChromaDB init failed: {e}")
        return None


def store_chunks_in_vector_db(
    chunks: List[str],
    material_id: int,
    subject_id: int,
) -> List[str]:
    """
    Embed and store text chunks in ChromaDB.
    Returns list of chroma_ids for each chunk.
    """
    collection = _get_chroma_collection()
    if collection is None:
        return [f"fallback-{i}" for i in range(len(chunks))]

    chroma_ids: List[str] = []
    for i, chunk in enumerate(chunks):
        cid = f"mat{material_id}_chunk{i}_{uuid.uuid4().hex[:8]}"
        try:
            collection.add(
                documents=[chunk],
                ids=[cid],
                metadatas=[{"material_id": material_id, "subject_id": subject_id, "chunk_index": i}],
            )
            chroma_ids.append(cid)
        except Exception as e:
            logger.warning(f"Failed to store chunk {i}: {e}")
            chroma_ids.append(f"error-{i}")

    return chroma_ids


def retrieve_relevant_chunks(
    query: str,
    subject_id: Optional[int] = None,
    material_id: Optional[int] = None,
    material_ids: Optional[List[int]] = None,
    n_results: int = 5,
    fallback_chunks: Optional[List[str]] = None,
) -> List[str]:
    """
    Query ChromaDB to retrieve semantically relevant chunks.
    If ChromaDB is unavailable or returns 0 results, uses keyword/similarity matching over fallback chunks.
    """
    collection = _get_chroma_collection()
    if collection is not None:
        try:
            # Scope the vector search by subject and (optionally) by material(s).
            # `material_id` scopes to a single file; `material_ids` scopes to a set
            # of files (NotebookLM-style source selection) via a ChromaDB `$in`.
            where = None
            if subject_id is not None or material_id is not None or (material_ids is not None and len(material_ids) > 0):
                where = {}
                if subject_id is not None:
                    where["subject_id"] = subject_id
                if material_id is not None:
                    where["material_id"] = material_id
                if material_ids is not None and len(material_ids) > 0:
                    where["material_id"] = {"$in": material_ids}
            results = collection.query(
                query_texts=[query],
                n_results=n_results,
                where=where,
            )
            docs = results.get("documents", [[]])[0]
            non_empty = [d for d in docs if d]
            if non_empty:
                return non_empty
        except Exception as e:
            logger.warning(f"ChromaDB query fallback: {e}")

    # Keyword/relevance fallback across provided chunks
    if fallback_chunks:
        q_words = set(query.lower().split())
        scored = []
        for c in fallback_chunks:
            c_words = set(c.lower().split())
            score = len(q_words.intersection(c_words))
            scored.append((score, c))
        scored.sort(key=lambda x: x[0], reverse=True)
        return [c for _, c in scored[:n_results]]

    return []


def delete_material_chunks(material_id: int) -> None:
    """Remove all vectors for a given material from ChromaDB."""
    collection = _get_chroma_collection()
    if collection is None:
        return
    try:
        collection.delete(where={"material_id": material_id})
    except Exception as e:
        logger.warning(f"Failed to delete chunks for material {material_id}: {e}")


def delete_subject_chunks(subject_id: int) -> None:
    """Remove all vectors for a given subject from ChromaDB (Slice 0.5).

    Without this, deleting a subject orphans its chunks (tagged with the dead
    subject_id), which global source-search would then incorrectly match.
    """
    collection = _get_chroma_collection()
    if collection is None:
        return
    try:
        collection.delete(where={"subject_id": subject_id})
    except Exception as e:
        logger.warning(f"Failed to delete chunks for subject {subject_id}: {e}")
